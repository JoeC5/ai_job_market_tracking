"""
generate_pptx.py
-----------------
Builds the weekly AI Job Market PowerPoint deck by:
  1. Copying the reference deck (4_22_2026_post.pptx) as the base
  2. Removing embedded charts and tables on data slides
  3. Inserting the PNG visuals produced by generate_visuals.py
  4. Updating the date in the title and slide headers

No LLM calls. No DB connection. Run generate_visuals.py first.

Usage:
    python generate_pptx.py

Output:
    outputs/AI_Job_Market_Report_YYYYMMDD.pptx

Requirements:
    pip install python-pptx
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Configuration ─────────────────────────────────────────────────────────────
# Update these two dates each week before running

REPORT_DATE = "05.19.2026"   # MM.DD.YYYY  — current report date
PREV_DATE   = "05.07.2026"   # MM.DD.YYYY  — previous report date

TEMPLATE_PATH = "ppt_template.pptx"
VISUALS_DIR   = "outputs/visuals"
OUTPUT_DIR    = "outputs"

# Visual filenames — produced by generate_visuals.py
VISUALS = {
    "global_trends":  "01_global_trends.png",
    "jobs_by_country":"02_jobs_by_country.png",
    "remote_trend":   "03_remote_hybrid_onsite.png",
    "industry_change":"04_industry_change.png",
    "tech_industries":"05_tech_industries.png",
    "non_tech":       "06_non_tech_industries.png",
    "role_change":    "07_role_type_change.png",
    "roles_by_count": "08_roles_by_count.png",
    "company_table":  "09_company_table.png",
    "us_map":         "10_us_state_map.png",
    "americas_map":   "11_americas_map.png",
    "europe_map":     "12_europe_map.png",
    "apac_map":       "13_apac_map.png",
}

# python-pptx shape type codes
CHART = 3
TABLE = 19

# ── Helpers ───────────────────────────────────────────────────────────────────

def vis(key):
    """Return the full path for a visual, with a warning if missing."""
    path = os.path.join(VISUALS_DIR, VISUALS[key])
    if not os.path.exists(path):
        print(f"   WARNING: visual not found — {path}")
    return path



def add_image(slide, img_path, left, top, width, height):
    """Add a PNG to the slide. All dimensions in inches."""
    if not os.path.exists(img_path):
        return None
    return slide.shapes.add_picture(
        img_path,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )


def update_title(slide, new_text):
    """Overwrite the title placeholder text, preserving font styling."""
    for shape in slide.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            tf = shape.text_frame
            lines = new_text.split("\n")
            for para_idx, para in enumerate(tf.paragraphs):
                line = lines[para_idx] if para_idx < len(lines) else ""
                if para.runs:
                    para.runs[0].text = line
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.add_run().text = line
            return


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_1_cover(slide):
    """Update the report date in the cover title."""
    for shape in slide.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if "as of" in run.text.lower():
                        run.text = f"as of {REPORT_DATE}"
                        return
    print("   NOTE: Cover date text not found — update manually if needed.")


def slide_2_methodology(slide):
    """Static methodology slide — no changes."""
    pass


def slide_3_global_trends(slide):
    """
    Template has two embedded charts (country bar top, global line bottom).
    Replace both with: global trend line (top) + remote/hybrid/onsite line (bottom).
    Original chart footprints:
        Chart 5 (top):    L=0.466  T=0.447  W=12.304  H=2.255
        Chart 1 (bottom): L=0.515  T=2.702  W=12.353  H=4.798
    """
  
    # Global trend line — fills the top chart slot and extends into the bottom
    add_image(slide, vis("global_trends"),
              left=0.47, top=0.45, width=12.35, height=4.30)

    # Remote / hybrid / onsite trend — bottom strip
    add_image(slide, vis("remote_trend"),
              left=0.47, top=4.85, width=12.35, height=2.55)


def slide_4_country_pct_change(slide):
    """
    Template has two tables: large country table (left) + remote summary table (right).
    Replace with: country bar chart (left) + remote/hybrid/onsite trend chart (right).
    Original table footprints:
        Table 2 (country): L=0.466  T=1.210  W=5.160  H=6.263
        Table 3 (remote):  L=7.344  T=1.508  W=4.086  H=1.019
    """
   
    update_title(slide,
        f'Global job data with "AI" as keyword by country\n')
        # f'Percentage change {PREV_DATE} to {REPORT_DATE}')

    # Jobs by country bar chart — left panel
    add_image(slide, vis("jobs_by_country"),
              left=0.47, top=1.50, width=12.35, height=6.26)

    # Remote / hybrid / onsite trend — right panel
    #add_image(slide, vis("remote_trend"),
    #          left=7.20, top=1.51, width=5.70, height=3.50)


def slide_5_industry(slide):
    """
    Template: tech bar chart (top-left), non-tech bar chart (bottom-left),
    pct-change table (right). Replace all with PNGs in matching footprints.
    Original footprints:
        Chart 6 (tech):    L=0.140  T=1.270  W=8.465  H=3.078
        Chart 7 (non-tech):L=-0.713 T=4.275  W=9.319  H=3.297  (left edge clips)
        Table 8 (pct chg): L=8.631  T=1.270  W=4.319  H=4.759
    """
    
    update_title(slide,
        f'Global job data with "AI" as keyword – By Industry\n'
        f'{PREV_DATE} to {REPORT_DATE}')

    # Tech industries — top left (Chart 6 footprint)
    add_image(slide, vis("tech_industries"),
              left=0.14, top=1.27, width=8.47, height=3.08)

    # Non-tech industries — bottom left (Chart 7 footprint, anchored at left edge)
    add_image(slide, vis("non_tech"),
              left=0.14, top=4.28, width=8.47, height=3.08)

    # Industry % change — right column (Table 8 footprint)
    add_image(slide, vis("industry_change"),
              left=8.63, top=1.27, width=4.57, height=6.05)


def slide_6_roles(slide):
    """
    Template: role % change chart (left), role count table (right).
    Replace with: roles by count (left), role % change (right).
    Original footprints:
        Chart 2 (role chg): L=0.512  T=1.143  W=7.998  H=6.071
        Table 3 (counts):   L=9.067  T=1.184  W=3.800  H=5.133
    """
    
    update_title(slide,
        f'Global job data – AI titled roles\n'
        f'{PREV_DATE} to {REPORT_DATE}')

    # Roles by count — left (Chart 2 footprint)
    add_image(slide, vis("roles_by_count"),
              left=0.51, top=1.14, width=8.00, height=6.07)

    # Role type % change — right (Table 3 footprint)
    add_image(slide, vis("role_change"),
              left=9.07, top=1.18, width=4.10, height=6.07)


def slide_7_keywords(slide):
    """Static keyword search reference table — no changes."""
    pass


def slide_8_companies(slide):
    """
    Template: highlights text box (left), company table (right).
    The text box stays; replace the table with the company table PNG.
    Original footprint:
        Table 3 (company): L=6.489  T=1.000  W=5.891  H=6.261
    """
    
    update_title(slide,
        f'Global job data with "AI" as keyword –\n'
        f'Company Openings: {PREV_DATE} to {REPORT_DATE}')

    # Company table PNG — right panel (matches original table footprint)
    add_image(slide, vis("company_table"),
              left=6.49, top=1.00, width=6.70, height=6.26)


def slide_9_us_map(slide):
    """
    Template: disclaimer text box (bottom-left), remote summary table (left),
    no embedded map (it was inserted manually in the original).
    Remove the table; insert the US state map PNG into the main content area.
    Original footprint reference:
        Table 6 (remote): L=0.210  T=2.327  W=3.027  H=0.806
    """
    
    update_title(slide,
        f'"AI" roles – United States\n{REPORT_DATE}')

    # US state map — right/center content area, leaves left panel for text
    add_image(slide, vis("us_map"),
              left=3.00, top=0.90, width=10.20, height=6.40)


def add_regional_map_slides(prs):
    """
    Append three new slides (Americas, Europe, APAC) using the Blank layout.
    Dark background matches the template's style on slides 3-9.
    """
    blank_layout = prs.slide_layouts[6]  # 'Blank'

    regions = [
        ("Americas — AI Job Openings by Country",   "americas_map"),
        ("Europe — AI Job Openings by Country",     "europe_map"),
        ("APAC & Middle East — AI Job Openings",    "apac_map"),
    ]

    for title_text, vis_key in regions:
        slide = prs.slides.add_slide(blank_layout)

        # Dark navy background
        bg = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(7.50)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        bg.line.fill.background()

        # Slide title text box
        txb = slide.shapes.add_textbox(
            Inches(0.47), Inches(0.08),
            Inches(12.35), Inches(0.65)
        )
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{title_text}   |   {REPORT_DATE}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Calibri"

        # Map image — full content area
        add_image(slide, vis(vis_key),
                  left=0.47, top=0.82, width=12.35, height=6.50)

        print(f"   Added: {title_text}")


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  AI Job Market — PowerPoint Builder")
    print(f"  Report date : {REPORT_DATE}")
    print(f"  Previous    : {PREV_DATE}")
    print("=" * 60)

    if not os.path.exists(TEMPLATE_PATH):
        sys.exit(f"\nERROR: Template not found → {TEMPLATE_PATH}\n")

    if not os.path.exists(VISUALS_DIR):
        sys.exit(
            f"\nERROR: Visuals directory not found → {VISUALS_DIR}\n"
            f"Run generate_visuals.py first.\n"
        )

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    print(f"\n→ Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)

    if len(slides) < 9:
        sys.exit(f"\nERROR: Expected 9 slides, found {len(slides)}\n")

    builders = [
        ("Cover",             slide_1_cover),
        ("Methodology",       slide_2_methodology),
        ("Global Trends",     slide_3_global_trends),
        ("Country % Change",  slide_4_country_pct_change),
        ("Industry",          slide_5_industry),
        ("AI Role Types",     slide_6_roles),
        ("Keyword Reference", slide_7_keywords),
        ("Company Table",     slide_8_companies),
        ("US State Map",      slide_9_us_map),
    ]

    for i, (label, fn) in enumerate(builders):
        print(f"\n→ Slide {i+1}: {label}")
        fn(slides[i])

    print("\n→ Adding regional map slides...")
    add_regional_map_slides(prs)

    date_tag = REPORT_DATE.replace(".", "")
    output_path = os.path.join(OUTPUT_DIR, f"AI_Job_Market_Report_{date_tag}.pptx")
    prs.save(output_path)

    print(f"\n{'=' * 60}")
    print(f"  Saved → {output_path}")
    print(f"  Total slides: {len(list(prs.slides))}")
    print("=" * 60)


if __name__ == "__main__":
    main()